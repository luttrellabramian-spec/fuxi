from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "docs" / "reports" / "伏羲_v0.2.5_项目汇报.pptx"

FONT = "Microsoft YaHei"

COLORS = {
    "paper": "F7F3EA",
    "ink": "111827",
    "muted": "6B7280",
    "line": "D9D0C3",
    "white": "FFFFFF",
    "red": "D94645",
    "teal": "0F766E",
    "gold": "C68A2E",
    "blue": "2563EB",
    "dark": "0E1117",
    "soft": "EFE7DA",
    "green": "15803D",
    "orange": "EA580C",
}


def rgb(hex_color: str) -> RGBColor:
    hex_color = hex_color.strip("#")
    return RGBColor(int(hex_color[:2], 16), int(hex_color[2:4], 16), int(hex_color[4:], 16))


def set_run(run, size=18, color="ink", bold=False):
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.color.rgb = rgb(COLORS[color])
    run.font.bold = bold


def fill(shape, color, transparency=0):
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(COLORS[color])
    if transparency:
        shape.fill.transparency = transparency
    shape.line.color.rgb = rgb(COLORS[color])


def no_line(shape):
    shape.line.fill.background()


def add_bg(slide, color="paper"):
    bg = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    fill(bg, color)
    no_line(bg)
    bg.z_order = 0
    return bg


def add_text(slide, x, y, w, h, text, size=20, color="ink", bold=False, align="left", line_spacing=1.08):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Pt(0)
    tf.margin_right = Pt(0)
    tf.margin_top = Pt(0)
    tf.margin_bottom = Pt(0)
    p = tf.paragraphs[0]
    p.alignment = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}[align]
    p.line_spacing = line_spacing
    run = p.add_run()
    run.text = text
    set_run(run, size=size, color=color, bold=bold)
    return box


def add_multiline(slide, x, y, w, h, lines, size=17, color="ink", bullet=False, gap=0):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Pt(0)
    tf.margin_right = Pt(0)
    tf.margin_top = Pt(0)
    tf.margin_bottom = Pt(0)
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = 0
        p.line_spacing = 1.18
        p.space_after = Pt(gap)
        if bullet:
            p.text = ""
            p._p.get_or_add_pPr().set("marL", "228600")
            p._p.get_or_add_pPr().set("indent", "-114300")
            run = p.add_run()
            run.text = f"• {line}"
        else:
            run = p.add_run()
            run.text = line
        set_run(run, size=size, color=color)
    return box


def add_title(slide, title, subtitle=None, index=None, dark=False):
    color = "white" if dark else "ink"
    muted = "soft" if dark else "muted"
    add_text(slide, 0.62, 0.42, 8.6, 0.48, title, size=26, color=color, bold=True)
    if subtitle:
        add_text(slide, 0.65, 0.96, 9.2, 0.32, subtitle, size=10.5, color=muted)
    if index is not None:
        add_text(slide, 11.95, 0.50, 0.8, 0.25, f"{index:02d}", size=10, color=muted, align="right")


def add_footer(slide, index, dark=False):
    color = "soft" if dark else "muted"
    add_text(slide, 0.62, 7.08, 3.6, 0.2, "伏羲 Fuxi v0.2.5 · AI Agent Engine", size=8.8, color=color)
    add_text(slide, 12.05, 7.08, 0.7, 0.2, f"{index:02d}", size=8.8, color=color, align="right")


def add_card(slide, x, y, w, h, title, body, accent="teal", title_size=15, body_size=11.5):
    card = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    fill(card, "white")
    card.line.color.rgb = rgb(COLORS["line"])
    card.adjustments[0] = 0.08
    stripe = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x), Inches(y), Inches(0.08), Inches(h))
    fill(stripe, accent)
    add_text(slide, x + 0.23, y + 0.2, w - 0.42, 0.3, title, size=title_size, color="ink", bold=True)
    add_text(slide, x + 0.23, y + 0.64, w - 0.42, h - 0.82, body, size=body_size, color="muted")
    return card


def add_badge(slide, x, y, text, color="teal"):
    box = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(1.68), Inches(0.36))
    fill(box, color)
    box.adjustments[0] = 0.45
    add_text(slide, x, y + 0.08, 1.68, 0.2, text, size=9.5, color="white", bold=True, align="center")


def add_arrow(slide, x1, y1, x2, y2, color="muted"):
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    line.line.color.rgb = rgb(COLORS[color])
    line.line.width = Pt(1.5)
    line.line.end_arrowhead = True
    return line


def add_box(slide, x, y, w, h, title, body="", color="white", accent="teal"):
    rect = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    fill(rect, color)
    rect.line.color.rgb = rgb(COLORS["line"])
    rect.adjustments[0] = 0.08
    add_text(slide, x + 0.15, y + 0.15, w - 0.3, 0.27, title, size=12.5, color="ink", bold=True, align="center")
    if body:
        add_text(slide, x + 0.18, y + 0.55, w - 0.36, h - 0.65, body, size=9.2, color="muted", align="center")
    accent_bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x + 0.12), Inches(y + h - 0.14), Inches(w - 0.24), Inches(0.04))
    fill(accent_bar, accent)
    return rect


def build_deck():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # 1. Cover
    slide = prs.slides.add_slide(blank)
    add_bg(slide, "dark")
    for i in range(8):
        shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.LINE_INVERSE, Inches(0.2 + i * 1.75), Inches(0.3), Inches(0.2 + i * 1.75), Inches(7.2))
        shape.line.color.rgb = rgb("1E293B")
        shape.line.transparency = 45
    fill(slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.68), Inches(1.08), Inches(0.08), Inches(4.8)), "red")
    add_text(slide, 0.92, 1.02, 5.0, 0.42, "FUXI / 伏羲", size=18, color="gold", bold=True)
    add_text(slide, 0.88, 1.6, 9.6, 1.0, "自进化 AI Agent 引擎", size=38, color="white", bold=True)
    add_text(slide, 0.92, 2.78, 8.3, 0.5, "从“会对话”走向“会记忆、会工具、会记录、会优化”的长期协作系统", size=17, color="soft")
    add_badge(slide, 0.92, 3.65, "v0.2.5 WIP", "teal")
    add_badge(slide, 2.78, 3.65, "Python + TypeScript", "blue")
    add_badge(slide, 4.82, 3.65, "gRPC + Memory", "gold")
    add_text(slide, 0.95, 6.66, 5.2, 0.25, "项目汇报 · 研究型工程原型", size=10.5, color="soft")

    # 2. Positioning
    slide = prs.slides.add_slide(blank)
    add_bg(slide)
    add_title(slide, "01 · 项目定位", "伏羲不是普通聊天机器人，而是一套 Agent 引擎验证场", 2)
    add_text(slide, 0.8, 1.65, 5.7, 0.72, "一句话：", size=16, color="red", bold=True)
    add_text(slide, 0.8, 2.12, 8.6, 1.1, "面向长期协作的自进化 AI Agent 引擎。", size=32, color="ink", bold=True)
    add_text(slide, 0.85, 3.45, 8.8, 0.72, "它的目标不是多回答一个问题，而是让 Agent 拥有持续工作所需的基础设施：记忆、工具、日志、反馈和迭代。", size=17, color="muted")
    add_card(slide, 0.82, 5.08, 3.55, 1.05, "研究价值", "验证 Agent 系统化能力，而非单点 prompt 技巧。", "teal")
    add_card(slide, 4.88, 5.08, 3.55, 1.05, "工程价值", "双运行时、gRPC、测试和启动脚本形成可维护工程。", "blue")
    add_card(slide, 8.94, 5.08, 3.55, 1.05, "作品价值", "可作为 AI 产品、设计作品和技术展示的底层证据。", "gold")
    add_footer(slide, 2)

    # 3. Problem
    slide = prs.slides.add_slide(blank)
    add_bg(slide)
    add_title(slide, "02 · 为什么要做", "当前 AI 使用常停留在一次性问答，难以形成长期工作流", 3)
    problems = [
        ("记不住", "上下文断裂，长期任务无法自然延续。", "red"),
        ("不可靠", "工具调用、外部 API、网络错误缺乏稳定兜底。", "orange"),
        ("不可追踪", "失败原因缺少结构化记录，难以复盘和优化。", "blue"),
        ("难进化", "系统没有从历史运行中学习的闭环。", "teal"),
    ]
    for i, (t, b, c) in enumerate(problems):
        add_card(slide, 0.85 + i * 3.06, 1.65, 2.55, 2.1, t, b, c, title_size=20, body_size=12.5)
    add_text(slide, 1.0, 4.55, 11.2, 0.72, "伏羲的核心判断：Agent 真正有用的地方，在于它能作为一个持续系统参与任务，而不是只做一次回答。", size=21, color="ink", bold=True, align="center")
    add_multiline(slide, 2.0, 5.65, 9.4, 0.75, ["因此 v0.2.5 重点补齐：性能边界、执行日志、工具可靠性、记忆检索、网关降级和演化入口。"], size=15, color="muted")
    add_footer(slide, 3)

    # 4. Goals
    slide = prs.slides.add_slide(blank)
    add_bg(slide)
    add_title(slide, "03 · 目标与验收", "v0.2 阶段围绕性能优化、自进化、已有能力修复展开", 4)
    add_card(slide, 0.82, 1.55, 3.65, 1.38, "性能优化", "热记忆 LRU、LLM 超时重试、gRPC 并发控制、温记忆检索性能。", "teal")
    add_card(slide, 4.82, 1.55, 3.65, 1.38, "自进化", "从结构化日志中分析失败模式，生成可解释、可审核的优化建议。", "blue")
    add_card(slide, 8.82, 1.55, 3.65, 1.38, "功能修正", "修复 v0.1 测试中暴露的问题，保持核心能力不退化。", "gold")
    metrics = [
        "热记忆：100 条 LRU / 单条 5000 char",
        "LLM：30 秒超时 / 3 次指数退避",
        "gRPC：并发上限 100 请求",
        "温记忆：FTS5 + 分页检索",
        "网关：25 秒超时保护 + 降级响应",
        "日志：JSONL 结构化记录，可用于后续分析",
    ]
    add_text(slide, 0.88, 3.62, 2.4, 0.3, "关键验收指标", size=18, color="ink", bold=True)
    add_multiline(slide, 0.95, 4.18, 11.2, 1.5, metrics, size=15.5, color="ink", bullet=True, gap=2)
    add_footer(slide, 4)

    # 5. Architecture diagram
    slide = prs.slides.add_slide(blank)
    add_bg(slide)
    add_title(slide, "04 · 总体架构", "TypeScript 管入口，Python 管智能核心，gRPC 负责跨运行时通信", 5)
    add_box(slide, 4.55, 1.28, 4.2, 0.8, "用户 / CLI / Web UI", "输入任务、查看结果、配置模型", "white", "red")
    add_arrow(slide, 6.65, 2.08, 6.65, 2.48, "muted")
    add_box(slide, 4.55, 2.48, 4.2, 0.92, "TypeScript Gateway", "HTTP / SSE / WebSocket / Settings UI", "white", "blue")
    add_arrow(slide, 6.65, 3.4, 6.65, 3.82, "muted")
    add_box(slide, 4.55, 3.82, 4.2, 0.72, "gRPC + Protocol Buffers", "跨语言协议边界", "white", "gold")
    add_arrow(slide, 6.65, 4.54, 6.65, 4.92, "muted")
    add_box(slide, 4.55, 4.92, 4.2, 0.95, "Python Fuxi Core", "ReAct / Tool Executor / LLM / Selector", "white", "teal")
    add_arrow(slide, 5.25, 5.88, 3.3, 6.25, "muted")
    add_arrow(slide, 6.65, 5.88, 6.65, 6.25, "muted")
    add_arrow(slide, 8.05, 5.88, 10.0, 6.25, "muted")
    add_box(slide, 1.45, 6.25, 2.25, 0.65, "热记忆", "当前会话", "white", "red")
    add_box(slide, 5.52, 6.25, 2.25, 0.65, "温记忆", "近期上下文", "white", "gold")
    add_box(slide, 9.58, 6.25, 2.25, 0.65, "冷记忆", "长期知识", "white", "teal")
    add_footer(slide, 5)

    # 6. Module matrix
    slide = prs.slides.add_slide(blank)
    add_bg(slide)
    add_title(slide, "05 · 核心模块", "伏羲的能力不是一个 prompt，而是一组可维护模块", 6)
    modules = [
        ("Gateway", "HTTP、SSE、WebSocket、设置页、限流和降级。", "blue"),
        ("Core Engine", "ReAct 主循环、会话上下文、工具调度。", "red"),
        ("LLM Client", "OpenAI 兼容 API、超时、重试、Circuit Breaker。", "gold"),
        ("Tool Executor", "参数校验、去重、缓存、并发和执行追踪。", "teal"),
        ("Memory", "Hot / Warm / Cold 三层记忆，兼顾速度与长期沉淀。", "orange"),
        ("Evolution", "查询分类、策略画像、工具排序、记忆优化和行为演化。", "green"),
    ]
    for i, item in enumerate(modules):
        x = 0.82 + (i % 3) * 4.05
        y = 1.55 + (i // 3) * 2.2
        add_card(slide, x, y, 3.55, 1.55, item[0], item[1], item[2])
    add_footer(slide, 6)

    # 7. Memory
    slide = prs.slides.add_slide(blank)
    add_bg(slide)
    add_title(slide, "06 · 三层记忆系统", "让 Agent 的上下文从短期响应走向长期协作", 7)
    add_card(slide, 0.9, 1.6, 3.3, 3.35, "热记忆 Hot", "面向当前会话。\n\nLRU 淘汰、大小上限、快速读取，避免无限增长。", "red", 18, 14)
    add_card(slide, 5.0, 1.6, 3.3, 3.35, "温记忆 Warm", "面向近期上下文。\n\nSQLite FTS5、BM25 排序、分页，适合搜索最近经验。", "gold", 18, 14)
    add_card(slide, 9.1, 1.6, 3.3, 3.35, "冷记忆 Cold", "面向长期知识。\n\n用于沉淀语义记忆和跨会话知识，为后续产品化预留空间。", "teal", 18, 14)
    add_arrow(slide, 4.25, 3.25, 4.85, 3.25, "muted")
    add_arrow(slide, 8.35, 3.25, 8.95, 3.25, "muted")
    add_text(slide, 1.1, 5.72, 11.1, 0.45, "设计重点：既要让当前对话足够快，也要让系统能逐步形成“可检索的经验”。", size=19, color="ink", bold=True, align="center")
    add_footer(slide, 7)

    # 8. Tool loop
    slide = prs.slides.add_slide(blank)
    add_bg(slide)
    add_title(slide, "07 · 工具执行闭环", "从“能调用工具”升级到“能可靠地管理工具”", 8)
    steps = [
        ("任务理解", "识别用户意图"),
        ("工具选择", "Selector 排序"),
        ("执行保护", "超时/重试/校验"),
        ("结果回写", "记忆与日志"),
        ("反馈优化", "成功率追踪"),
    ]
    xs = [0.8, 3.1, 5.4, 7.7, 10.0]
    for i, (title, body) in enumerate(steps):
        add_box(slide, xs[i], 2.15, 1.85, 1.18, title, body, "white", ["red", "gold", "teal", "blue", "green"][i])
        if i < len(steps) - 1:
            add_arrow(slide, xs[i] + 1.86, 2.73, xs[i + 1] - 0.08, 2.73, "muted")
    add_multiline(slide, 1.2, 4.55, 10.7, 1.2, [
        "工具层不是简单调用函数，而是把每次执行都纳入可观察、可复盘、可降权的系统。",
        "这为后续“哪类工具容易失败、哪种策略更有效”提供了数据基础。",
    ], size=18, color="ink")
    add_footer(slide, 8)

    # 9. Evolution
    slide = prs.slides.add_slide(blank)
    add_bg(slide)
    add_title(slide, "08 · 自进化机制", "自动分析，但不自动越权修改核心系统", 9)
    add_box(slide, 0.9, 1.8, 2.2, 1.0, "结构化日志", "JSONL 记录每轮执行", "white", "blue")
    add_arrow(slide, 3.15, 2.3, 4.0, 2.3, "muted")
    add_box(slide, 4.05, 1.8, 2.2, 1.0, "失败模式分析", "timeout / tool_failure", "white", "red")
    add_arrow(slide, 6.3, 2.3, 7.15, 2.3, "muted")
    add_box(slide, 7.2, 1.8, 2.2, 1.0, "优化建议", "策略、工具、记忆", "white", "gold")
    add_arrow(slide, 9.45, 2.3, 10.3, 2.3, "muted")
    add_box(slide, 10.35, 1.8, 2.2, 1.0, "人工确认", "可解释、可回滚", "white", "teal")
    add_card(slide, 1.05, 4.2, 5.2, 1.35, "边界原则", "系统可以生成建议，但生产行为需要人工审核。进化优化能力，不改变核心定位。", "red")
    add_card(slide, 7.0, 4.2, 5.2, 1.35, "设计意义", "避免 Agent 自行改写自身目标，同时保留持续变强的工程入口。", "teal")
    add_footer(slide, 9)

    # 10. Validation
    slide = prs.slides.add_slide(blank)
    add_bg(slide)
    add_title(slide, "09 · 当前验证结果", "v0.2.5 已形成一条可重复检查的核心路径", 10)
    add_text(slide, 0.95, 1.55, 3.8, 0.42, "已通过", size=18, color="green", bold=True)
    add_card(slide, 0.95, 2.05, 3.55, 1.45, "TypeScript Gateway", "`npm run build` 通过，网关代码可编译。", "blue")
    add_card(slide, 4.9, 2.05, 3.55, 1.45, "Python 核心测试", "热记忆、温记忆、工具注册：115 passed。", "green")
    add_card(slide, 8.85, 2.05, 3.55, 1.45, "核心检查脚本", "`scripts/check_core.ps1` 一键验证基础路径。", "teal")
    add_text(slide, 0.95, 4.38, 3.8, 0.35, "版本与提交", size=18, color="ink", bold=True)
    add_multiline(slide, 1.0, 4.95, 10.8, 0.85, [
        "当前版本：v0.2.5 WIP",
        "最新提交：cbeb52c chore: prepare Fuxi v0.2.5 workspace",
        "项目根目录：fuxi_v0.2.5",
    ], size=15.5, color="ink", bullet=True)
    add_footer(slide, 10)

    # 11. Demo
    slide = prs.slides.add_slide(blank)
    add_bg(slide)
    add_title(slide, "10 · 演示路径", "把工程能力变成能被别人看懂的操作流程", 11)
    add_card(slide, 0.9, 1.55, 3.4, 1.35, "1. 启动", "Windows 根目录运行：\n.\\start.bat", "red")
    add_card(slide, 4.95, 1.55, 3.4, 1.35, "2. 访问", "http://localhost:18789/chat/ui\n/settings/ui\n/health", "blue")
    add_card(slide, 9.0, 1.55, 3.4, 1.35, "3. 展示", "一次对话、一次工具调用、一次记忆写入、一次日志查看。", "teal")
    add_text(slide, 0.95, 4.08, 3.3, 0.35, "汇报 Demo 建议", size=18, color="ink", bold=True)
    add_multiline(slide, 1.0, 4.62, 10.9, 1.2, [
        "不要从代码开始讲，先让观众看到“输入任务 → Agent 调用工具 → 写入记忆 → 可追踪日志”。",
        "汇报时强调：伏羲的价值在于系统闭环，而不只是一次回答的质量。",
    ], size=16.5, color="muted")
    add_footer(slide, 11)

    # 12. Limitations
    slide = prs.slides.add_slide(blank)
    add_bg(slide)
    add_title(slide, "11 · 当前限制", "原型已成型，但离可发布产品还有几处关键距离", 12)
    items = [
        ("安全", "认证、鉴权、工具权限和公网暴露策略仍需加强。", "red"),
        ("前端", "网关内嵌 UI 较多，后续应拆成独立产品界面。", "orange"),
        ("测试", "完整测试包含真实 API / MCP / 安全场景，需要分层标记。", "blue"),
        ("交付", "需要固定演示数据、启动脚本和 Demo 流程。", "teal"),
    ]
    for i, (t, b, c) in enumerate(items):
        add_card(slide, 0.9 + (i % 2) * 6.1, 1.7 + (i // 2) * 2.05, 5.35, 1.35, t, b, c, 17, 13)
    add_footer(slide, 12)

    # 13. Roadmap
    slide = prs.slides.add_slide(blank)
    add_bg(slide)
    add_title(slide, "12 · 下一步", "从工程原型，走向可展示产品与作品集核心案例", 13)
    roadmap = [
        ("绿色演示路径", "start.bat → /health → /chat/ui → 工具 → 记忆 → 日志", "red"),
        ("UI 产品化", "拆出独立前端，让伏羲有一个真正可展示的界面", "blue"),
        ("测试分层", "core / integration / real_api / security 分层运行", "gold"),
        ("产品外壳", "以“念”或“造物局 AI”承接伏羲底层能力", "teal"),
    ]
    for i, (t, b, c) in enumerate(roadmap):
        add_card(slide, 0.9, 1.45 + i * 1.28, 11.4, 0.9, t, b, c, 15, 12)
    add_text(slide, 1.0, 6.55, 11.2, 0.28, "结论：伏羲已经不是想法，而是可以继续产品化、作品集化、Demo 化的底层系统。", size=17, color="ink", bold=True, align="center")
    add_footer(slide, 13)

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    return OUT


if __name__ == "__main__":
    path = build_deck()
    print(path)
